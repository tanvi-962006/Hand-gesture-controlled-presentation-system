import os
import uuid
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
from flask import Flask, render_template, request, jsonify, url_for
from werkzeug.utils import secure_filename

# Import our Python Computer Vision & Shape Recognition Pipelines
from gesture_pipeline import HandGestureClassifier, PresentationAnalyticsTracker
from shape_recognizer import GeometricShapeRecognizer

app = Flask(__name__)
app.config['SECRET_KEY'] = 'aerosense_ai_spatial_studio_key_2026'
app.config['UPLOAD_FOLDER'] = os.path.join(app.root_path, 'static', 'uploads')
app.config['SLIDES_FOLDER'] = os.path.join(app.root_path, 'static', 'slides')
app.config['DEMO_FOLDER'] = os.path.join(app.config['SLIDES_FOLDER'], 'demo')
app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024  # 32 MB limit

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['SLIDES_FOLDER'], exist_ok=True)
os.makedirs(app.config['DEMO_FOLDER'], exist_ok=True)

@app.after_request
def add_header(response):
    response.headers['Cache-Control'] = 'no-store, no-cache, must-revalidate, post-check=0, pre-check=0, max-age=0'
    response.headers['Pragma'] = 'no-cache'
    response.headers['Expires'] = '-1'
    return response

@app.route('/favicon.ico')
def favicon():
    from flask import Response
    svg = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
        <defs>
            <linearGradient id="g" x1="0%" y1="0%" x2="100%" y2="100%">
                <stop offset="0%" stop-color="#00f0ff"/>
                <stop offset="100%" stop-color="#8b5cf6"/>
            </linearGradient>
        </defs>
        <circle cx="50" cy="50" r="46" fill="#0c0c1e" stroke="url(#g)" stroke-width="4"/>
        <text x="50" y="65" font-size="44" text-anchor="middle" fill="#00f0ff" font-family="system-ui, -apple-system, sans-serif">✨</text>
    </svg>'''
    return Response(svg, mimetype='image/svg+xml')

# Global Active Presentation State & Analytics Tracker
analytics_tracker = PresentationAnalyticsTracker()

active_presentation = {
    "id": "demo",
    "title": "AeroSense AI — Spatial Computing & Presentation Deck",
    "filename": "demo_slides",
    "slides": [],
    "notes": [],
    "current_index": 0,
    "total_slides": 0
}

# Demo slides metadata with speaker notes
DEMO_SLIDES_METADATA = [
    {
        "tag": "SPATIAL COMPUTING",
        "title": "AeroSense AI: Spatial Computing Studio",
        "subtitle": "Touchless Air-Canvas, Smart Shape Recognition & 3D Spatial Inspector",
        "cards": [
            ("Air Whiteboard Canvas", "Draw notes, equations, and diagrams freely in mid-air"),
            ("Geometric Shape Snapping", "Rough air circles and boxes snap into perfect vector geometry"),
            ("3D Spatial Object Viewer", "Inspect and rotate 3D scientific structures with hand orbit controls"),
            ("Multi-Modal Deck Presenter", "Voice and gesture-controlled presentations with live teleprompter")
        ],
        "footer": "Slide 1 / 5  •  AeroSense AI Spatial Studio",
        "speaker_notes": "Welcome to AeroSense AI! This is a complete Spatial Computing platform featuring an infinite Air Whiteboard, automatic geometric shape snapping, and interactive 3D model inspection."
    },
    {
        "tag": "GESTURE CHEATSHEET",
        "title": "Natural Hand Gestures Guide",
        "subtitle": "Intuitive AI-recognized gestures for presentation & spatial drawing",
        "cards": [
            ("Swipe / Point Right", "Advances presentation to the next slide smoothly"),
            ("Swipe / Point Left", "Returns presentation to the previous slide"),
            ("Pinch Thumb & Index", "Enables air drawing pen on whiteboard or presentation"),
            ("Closed Fist", "Activates spatial Air Eraser to wipe notes on screen"),
            ("Victory Sign (V)", "Triggers celebration confetti and snaps active hand drawings")
        ],
        "footer": "Slide 2 / 5  •  Spatial Gesture Dictionary",
        "speaker_notes": "Here are our core spatial hand gestures. Swiping navigates slides, pinching thumb and index lets you write in the air, and a closed fist acts as an air eraser."
    },
    {
        "tag": "COMPUTER VISION",
        "title": "Pure Geometric AI Architecture",
        "subtitle": "Real-time landmark kinematics and Euclidean variance classification",
        "cards": [
            ("Sub-25ms Latency", "Runs at 60 FPS via client WebAssembly and Python REST classifier"),
            ("Shape Variance Mathematics", "Analyzes radial variance to detect circles, rectangles, and lines"),
            ("Vector Document Pipeline", "150 DPI page rendering using PyMuPDF and python-pptx extraction"),
            ("Zero Transformer Overhead", "Lightweight, interpretable mathematical algorithms")
        ],
        "footer": "Slide 3 / 5  •  Tech Stack & Computer Vision",
        "speaker_notes": "Under the hood, we use pure geometric mathematics: Euclidean distance variance, corner angle detection, and dynamic velocity vectors without heavy black-box transformer overhead."
    },
    {
        "tag": "3D INTERACTION",
        "title": "Interactive 3D Diagram Inspector",
        "subtitle": "Two-hand spatial rotation and multi-axis orbit manipulation",
        "cards": [
            ("3D Geometric Solids", "Inspect icosahedrons, crystalline lattices, and complex polyhedra"),
            ("Biomolecular Visualizer", "Examine molecular structures and double-helix DNA chains"),
            ("Two-Hand Distance Zoom", "Separate or bring hands together to zoom into diagrams in 3D"),
            ("Air-Grab Rotation", "Open palm spatial movement rotates 3D objects in real time")
        ],
        "footer": "Slide 4 / 5  •  3D Spatial Inspection",
        "speaker_notes": "Our 3D Spatial Inspector allows presenters to load interactive 3D models and rotate them in mid-air using open palm movements."
    },
    {
        "tag": "GET STARTED",
        "title": "Ready to Experience AeroSense AI",
        "subtitle": "Switch between Air Whiteboard, Spatial Deck, and 3D Model Inspector",
        "cards": [
            ("Switch Modes", "Use the top floating pill to jump between Whiteboard, Deck, and 3D Inspector"),
            ("Upload Custom Decks", "Upload your own .pdf or .pptx presentation files anytime"),
            ("Voice AI Commands", "Speak 'Next', 'Back', 'Clear', 'Spotlight' for dual-modal control"),
            ("Export Artifacts", "Save annotated slides and air-drawings directly to your computer")
        ],
        "footer": "Slide 5 / 5  •  Ready to Present",
        "speaker_notes": "You are now ready to explore all 3 modes! Switch to the Air Whiteboard to test shape snapping, or open the 3D Inspector to rotate 3D structures."
    }
]

def generate_demo_slides():
    """Generates crisp 1920x1080 demo slides with cards and typography."""
    demo_dir = app.config['DEMO_FOLDER']
    generated_urls = []
    generated_notes = []
    width, height = 1920, 1080

    for idx, slide_data in enumerate(DEMO_SLIDES_METADATA, 1):
        filename = f"demo_slide_{idx}.png"
        filepath = os.path.join(demo_dir, filename)
        rel_url = f"/static/slides/demo/{filename}"
        generated_urls.append(rel_url)
        generated_notes.append(slide_data.get("speaker_notes", ""))

        img = Image.new("RGB", (width, height), "#06060f")
        draw = ImageDraw.Draw(img)

        # Smooth vertical gradient background
        for y in range(height):
            ratio = y / height
            r = int(6 + ratio * 10)
            g = int(6 + ratio * 14)
            b = int(18 + ratio * 32)
            draw.line([(0, y), (width, y)], fill=(r, g, b))

        # Ambient header glow line
        draw.rectangle([(0, 0), (width, 8)], fill="#8b5cf6")
        draw.rectangle([(0, 8), (600, 12)], fill="#00bbf9")

        # Fonts
        try:
            tag_font = ImageFont.truetype("arialbd.ttf", 22)
            title_font = ImageFont.truetype("arialbd.ttf", 52)
            sub_font = ImageFont.truetype("arial.ttf", 28)
            card_title_font = ImageFont.truetype("arialbd.ttf", 28)
            card_desc_font = ImageFont.truetype("arial.ttf", 22)
            footer_font = ImageFont.truetype("arial.ttf", 20)
        except Exception:
            tag_font = ImageFont.load_default()
            title_font = ImageFont.load_default()
            sub_font = ImageFont.load_default()
            card_title_font = ImageFont.load_default()
            card_desc_font = ImageFont.load_default()
            footer_font = ImageFont.load_default()

        # Category Tag Pill
        tag_text = slide_data["tag"]
        draw.rounded_rectangle([(100, 70), (360, 106)], radius=18, fill=(139, 92, 246, 40), outline="#8b5cf6", width=2)
        draw.text((120, 78), tag_text, fill="#c4b5fd", font=tag_font)

        # Slide Main Title & Subtitle
        draw.text((100, 130), slide_data["title"], fill="#ffffff", font=title_font)
        draw.text((100, 205), slide_data["subtitle"], fill="#94a3b8", font=sub_font)

        # Cards Section
        cards = slide_data["cards"]
        start_y = 280
        
        if len(cards) <= 4:
            positions = [
                (100, start_y, 920, start_y + 300),
                (980, start_y, 1800, start_y + 300),
                (100, start_y + 340, 920, start_y + 640),
                (980, start_y + 340, 1800, start_y + 640)
            ]
            for i, (ctitle, cdesc) in enumerate(cards):
                if i >= len(positions):
                    break
                x1, y1, x2, y2 = positions[i]
                draw.rounded_rectangle([(x1, y1), (x2, y2)], radius=20, fill=(16, 16, 32), outline=(50, 50, 90), width=2)
                draw.rounded_rectangle([(x1 + 30, y1 + 32), (x1 + 36, y1 + 72)], radius=3, fill="#8b5cf6")
                draw.text((x1 + 55, y1 + 35), ctitle, fill="#ffffff", font=card_title_font)
                
                words = cdesc.split()
                line = ""
                y_text = y1 + 105
                for w in words:
                    test_line = line + " " + w if line else w
                    if len(test_line) > 42:
                        draw.text((x1 + 35, y_text), line, fill="#cbd5e1", font=card_desc_font)
                        y_text += 38
                        line = w
                    else:
                        line = test_line
                if line:
                    draw.text((x1 + 35, y_text), line, fill="#cbd5e1", font=card_desc_font)

        # Footer Bar
        draw.line([(100, 990), (1800, 990)], fill=(40, 40, 70), width=1)
        draw.text((100, 1015), slide_data["footer"], fill="#64748b", font=footer_font)
        draw.text((1500, 1015), "AeroSense AI Spatial Studio", fill="#8b5cf6", font=footer_font)

        img.save(filepath, "PNG", quality=95)

    return generated_urls, generated_notes

demo_urls, demo_notes = generate_demo_slides()
active_presentation["slides"] = demo_urls
active_presentation["notes"] = demo_notes
active_presentation["total_slides"] = len(demo_urls)

def process_pdf(pdf_path, presentation_id):
    """Converts PDF pages into slide images and extracts text for teleprompter."""
    out_dir = os.path.join(app.config['SLIDES_FOLDER'], presentation_id)
    os.makedirs(out_dir, exist_ok=True)

    doc = fitz.open(pdf_path)
    slide_urls = []
    slide_notes = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=150)
        filename = f"slide_{page_num + 1}.png"
        filepath = os.path.join(out_dir, filename)
        pix.save(filepath)
        slide_urls.append(f"/static/slides/{presentation_id}/{filename}")

        raw_text = page.get_text().strip()
        if raw_text:
            lines = [line.strip() for line in raw_text.split('\n') if line.strip()]
            summary_notes = " • ".join(lines[:5]) if lines else "Slide text content."
            slide_notes.append(summary_notes)
        else:
            slide_notes.append(f"Notes for Slide {page_num + 1}")

    doc.close()
    return slide_urls, slide_notes

def process_pptx(pptx_path, presentation_id):
    """Parses PPTX presentation content, renders slide images, and extracts speaker notes."""
    out_dir = os.path.join(app.config['SLIDES_FOLDER'], presentation_id)
    os.makedirs(out_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    slide_urls = []
    slide_notes = []
    width, height = 1920, 1080

    try:
        title_font = ImageFont.truetype("arialbd.ttf", 48)
        text_font = ImageFont.truetype("arial.ttf", 26)
        footer_font = ImageFont.truetype("arial.ttf", 20)
    except Exception:
        title_font = ImageFont.load_default()
        text_font = ImageFont.load_default()
        footer_font = ImageFont.load_default()

    for idx, slide in enumerate(prs.slides, 1):
        filename = f"slide_{idx}.png"
        filepath = os.path.join(out_dir, filename)

        slide_title = f"Slide {idx}"
        slide_texts = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    if shape == slide.shapes[0] and idx == 1:
                        slide_title = text
                    elif hasattr(shape, "is_placeholder") and shape.is_placeholder and "Title" in str(shape.placeholder_format.type):
                        slide_title = text
                    else:
                        slide_texts.append(text)

        speaker_note = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                speaker_note = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        if not speaker_note and slide_texts:
            speaker_note = " • ".join(slide_texts[:4])
        elif not speaker_note:
            speaker_note = f"Key topics and discussion points for Slide {idx}."

        slide_notes.append(speaker_note)

        img = Image.new("RGB", (width, height), "#06060f")
        draw = ImageDraw.Draw(img)

        for y in range(height):
            ratio = y / height
            draw.line([(0, y), (width, y)], fill=(int(6 + ratio * 10), int(6 + ratio * 14), int(18 + ratio * 32)))

        draw.rectangle([(0, 0), (width, 8)], fill="#8b5cf6")
        draw.rounded_rectangle([(80, 50), (1840, 160)], radius=20, fill=(16, 16, 32), outline=(50, 50, 90), width=2)
        draw.text((120, 80), slide_title[:60], fill="#ffffff", font=title_font)

        draw.rounded_rectangle([(80, 190), (1840, 960)], radius=20, fill=(14, 14, 28), outline=(45, 45, 80), width=2)

        y_offset = 240
        if slide_texts:
            for text_line in slide_texts[:10]:
                draw.ellipse([(120, y_offset + 8), (134, y_offset + 22)], fill="#8b5cf6")
                draw.text((160, y_offset), text_line[:90], fill="#f1f5f9", font=text_font)
                y_offset += 65
        else:
            draw.text((160, 400), "PowerPoint Slide Content", fill="#64748b", font=text_font)

        draw.text((100, 1000), f"Slide {idx} of {len(prs.slides)}", fill="#64748b", font=footer_font)
        draw.text((1500, 1000), "AeroSense AI Spatial Studio", fill="#8b5cf6", font=footer_font)

        img.save(filepath, "PNG", quality=95)
        slide_urls.append(f"/static/slides/{presentation_id}/{filename}")

    return slide_urls, slide_notes


# ── REST API ROUTES ──

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/slides', methods=['GET'])
def get_slides():
    """Returns active presentation state and slide list."""
    return jsonify({
        "status": "success",
        "presentation": active_presentation
    })

@app.route('/api/notes', methods=['GET'])
def get_notes():
    """Returns speaker notes for the current active slide."""
    curr_idx = active_presentation["current_index"]
    notes = active_presentation["notes"]
    note_text = notes[curr_idx] if (notes and 0 <= curr_idx < len(notes)) else "No speaker notes recorded for this slide."
    return jsonify({
        "status": "success",
        "current_index": curr_idx,
        "note": note_text
    })

@app.route('/api/navigate', methods=['POST'])
def navigate_slide():
    """Navigates slide index and updates analytics."""
    data = request.json or {}
    action = data.get('action')
    target_idx = data.get('index', None)

    total = active_presentation["total_slides"]
    current = active_presentation["current_index"]
    new_idx = current

    if action == 'next':
        if current < total - 1:
            new_idx = current + 1
    elif action == 'prev':
        if current > 0:
            new_idx = current - 1
    elif action == 'goto' and target_idx is not None:
        if 0 <= target_idx < total:
            new_idx = target_idx

    if new_idx != current:
        active_presentation["current_index"] = new_idx
        analytics_tracker.record_slide_change(new_idx)

    return jsonify({
        "status": "success",
        "presentation": active_presentation
    })

@app.route('/api/recognize_shape', methods=['POST'])
def recognize_shape():
    """
    Analyzes drawn 2D coordinates from Air Whiteboard and returns
    recognized geometric primitive (Circle, Rectangle, Triangle, Line, Freehand).
    """
    data = request.json or {}
    points = data.get('points', [])
    shape_result = GeometricShapeRecognizer.recognize(points)
    return jsonify({
        "status": "success",
        "shape": shape_result
    })

@app.route('/api/analytics', methods=['GET'])
def get_analytics():
    """Returns live presentation & workspace analytics."""
    summary = analytics_tracker.get_summary(active_presentation["total_slides"])
    return jsonify({
        "status": "success",
        "analytics": summary
    })

@app.route('/api/analytics/log_gesture', methods=['POST'])
def log_gesture():
    """Logs recognized gesture event into tracker."""
    data = request.json or {}
    gesture = data.get('gesture')
    if gesture:
        analytics_tracker.record_gesture(gesture)
    return jsonify({"status": "success"})

@app.route('/api/analytics/reset', methods=['POST'])
def reset_analytics():
    analytics_tracker.reset()
    return jsonify({"status": "success", "message": "Analytics reset."})

@app.route('/api/process_frame', methods=['POST'])
def process_frame():
    """Python Server-Side Euclidean Landmark Classifier endpoint."""
    data = request.json or {}
    landmarks = data.get('landmarks', [])
    result = HandGestureClassifier.classify(landmarks)
    return jsonify({
        "status": "success",
        "result": result
    })

def parse_voice_action(text: str):
    if not text:
        return None, None
    t = text.lower().strip()
    
    # 1. Next Slide (Phonetic variations: next, nex, nest, necks, text, max, aage, agla, chalo, etc.)
    if any(w in t for w in ['next', 'nex', 'nest', 'necks', 'nxt', 'aage', 'agla', 'forward', 'right', 'chalo', 'next page', 'go next']):
        return 'next_slide', 'Next Slide'
    
    # 2. Previous Slide (Phonetic variations: back, bac, pack, beck, piche, peeche, previous, prev, wapas, etc.)
    if any(w in t for w in ['back', 'bac', 'pack', 'beck', 'piche', 'peeche', 'previous', 'prev', 'left', 'wapas', 'pichla', 'go back']):
        return 'prev_slide', 'Previous Slide'
    
    # 3. Whiteboard (Phonetic: whiteboard, white board, board, bord, draw, drawing)
    if any(w in t for w in ['whiteboard', 'white board', 'board', 'bord', 'draw mode', 'drawing', 'canvas', 'likho']):
        return 'whiteboard', 'Switched to Air Whiteboard'
    
    # 4. Presentation Deck
    if any(w in t for w in ['presentation', 'slide', 'slides', 'deck', 'ppt']):
        return 'deck', 'Switched to Spatial Deck'
    
    # 5. 3D Model
    if any(w in t for w in ['3d', 'three d', '3-d', 'three-d', 'model', 'inspector', 'hologram', 'cube', 'spatial']):
        return 'model3d', 'Switched to 3D Spatial Inspector'
    
    # 6. Laser Pointer (Phonetic: laser, lazer, leser, layer, lezer, lejar, dot, point)
    if any(w in t for w in ['laser', 'lazer', 'leser', 'lezer', 'lejar', 'pointer', 'point', 'red dot', 'red laser']):
        return 'laser', 'Laser Pointer Active'
    
    # 7. Creative Tools
    if any(w in t for w in ['rainbow', 'color', 'colour']):
        return 'rainbow', 'Rainbow Pen Active'
    if any(w in t for w in ['sparkler', 'sparkle', 'magic']):
        return 'sparkler', 'Sparkler Brush Active'
    if any(w in t for w in ['shape', 'geometry', 'circle', 'square', 'triangle']):
        return 'shape', 'Shape AI Active'
    if any(w in t for w in ['spotlight', 'torch', 'focus', 'beam']):
        return 'spotlight', 'Spotlight Beam Active'
    if any(w in t for w in ['eraser', 'rubber', 'mitao', 'erase']):
        return 'eraser', 'Air Eraser Active'
    if any(w in t for w in ['pen', 'pencil', 'marker', 'write']):
        return 'pen', 'Neon Pen Active'
    
    # 8. Actions
    if any(w in t for w in ['clear', 'clean', 'cler', 'clea', 'saaf', 'saf', 'delete', 'remove all']):
        return 'clear', 'Canvas Cleared'
    if any(w in t for w in ['camera on', 'start camera', 'turn on camera', 'open camera']):
        return 'camera_on', 'Camera Started'
    if any(w in t for w in ['camera off', 'stop camera', 'turn off camera', 'close camera']):
        return 'camera_off', 'Camera Stopped'
    if any(w in t for w in ['spin', 'rotate', 'ghumo', 'auto spin']):
        return 'spin', 'Auto Spin Toggled'
    if any(w in t for w in ['wireframe', 'mesh', 'skeleton']):
        return 'wireframe', 'Wireframe Toggled'
    if any(w in t for w in ['center', 'reset 3d', 'reset model']):
        return 'center', '3D Model Centered'
    if any(w in t for w in ['fullscreen', 'full screen', 'maximize']):
        return 'fullscreen', 'Toggled Fullscreen'
    return None, None

@app.route('/api/voice_command', methods=['POST'])
def process_voice_command():
    """Dual-Engine Voice Command Processor (Direct Text / Audio Stream)."""
    import speech_recognition as sr
    transcript = ""

    # 1. Direct JSON text payload
    if request.is_json:
        data = request.json or {}
        transcript = data.get('text', '').strip()
    
    # 2. Audio file upload
    if not transcript and 'audio' in request.files:
        audio_file = request.files['audio']
        try:
            recognizer = sr.Recognizer()
            recognizer.dynamic_energy_threshold = True
            recognizer.energy_threshold = 100  # Extra sensitive for distance speech
            recognizer.pause_threshold = 0.4
            with sr.AudioFile(audio_file) as source:
                audio_data = recognizer.record(source)
                try:
                    # Recognize with Indian English / Hinglish support
                    transcript = recognizer.recognize_google(audio_data, language='en-IN')
                except Exception:
                    # Fallback to general English
                    transcript = recognizer.recognize_google(audio_data, language='en-US')
        except Exception as e:
            print(f"Server Speech Recognition note: {e}")

    if not transcript:
        return jsonify({"status": "no_speech", "message": "No command heard"}), 200

    action, label = parse_voice_action(transcript)
    return jsonify({
        "status": "success",
        "transcript": transcript,
        "action": action,
        "message": label or f'Recognized: "{transcript}"'
    })

@app.route('/api/upload', methods=['POST'])
def upload_presentation():
    """Uploads PDF or PPTX file and converts to vector slides with notes."""
    if 'file' not in request.files:
        return jsonify({"status": "error", "message": "No file uploaded"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"status": "error", "message": "No selected file"}), 400

    filename = secure_filename(file.filename)
    ext = os.path.splitext(filename)[1].lower()

    if ext not in ['.pdf', '.pptx', '.ppt']:
        return jsonify({"status": "error", "message": "Unsupported file format. Please upload PDF or PPTX."}), 400

    presentation_id = str(uuid.uuid4())[:8]
    saved_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{presentation_id}_{filename}")
    file.save(saved_path)

    try:
        if ext == '.pdf':
            slides, notes = process_pdf(saved_path, presentation_id)
        else:
            slides, notes = process_pptx(saved_path, presentation_id)

        if not slides:
            return jsonify({"status": "error", "message": "Could not extract slides from file."}), 400

        active_presentation["id"] = presentation_id
        active_presentation["title"] = os.path.splitext(filename)[0].replace('_', ' ').title()
        active_presentation["filename"] = filename
        active_presentation["slides"] = slides
        active_presentation["notes"] = notes
        active_presentation["current_index"] = 0
        active_presentation["total_slides"] = len(slides)

        analytics_tracker.reset()

        return jsonify({
            "status": "success",
            "message": f"Successfully processed {len(slides)} slides with notes!",
            "presentation": active_presentation
        })
    except Exception as e:
        print(f"Error processing upload: {e}")
        return jsonify({"status": "error", "message": f"Failed to process file: {str(e)}"}), 500

@app.route('/api/reset_demo', methods=['POST'])
def reset_demo():
    """Resets presentation to default demo deck."""
    active_presentation["id"] = "demo"
    active_presentation["title"] = "AeroSense AI — Spatial Computing & Presentation Deck"
    active_presentation["filename"] = "demo_slides"
    active_presentation["slides"] = demo_urls
    active_presentation["notes"] = demo_notes
    active_presentation["current_index"] = 0
    active_presentation["total_slides"] = len(demo_urls)

    analytics_tracker.reset()

    return jsonify({
        "status": "success",
        "message": "Reset to default demo presentation.",
        "presentation": active_presentation
    })

if __name__ == '__main__':
    print("Starting AeroSense AI Spatial Studio on http://127.0.0.1:5000")
    app.run(host='0.0.0.0', port=5000, debug=True)
